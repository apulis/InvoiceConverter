#!/usr/bin/env Python
# coding=utf-8

import sys, datetime, fitz, os, codecs
from pptx import Presentation
from pptx.util import Cm
import comtypes.client
from PyPDF2 import PdfFileReader, PdfFileWriter, PdfFileMerger

gdict = {0:u'零',1:u'壹',2:u'贰',3:u'叁',4:u'肆',5:u'伍',6:u'陆',7:u'柒',8:u'捌',9:u'玖',10:u'拾'}

def numToCN(num):
    cn = ''

    if num < 10:
        cn = gdict[num]
    elif num > 10 and num < 100 :
        if num % 10 == 0:
            cn = '{0}拾'.format(gdict[num // 10])
        else:
            cn = '{0}拾{1}'.format(gdict[num // 10], gdict[num % 10])
    return cn

class InvoiceConverter(object):
    def __init__(self, name, totalPage = None, totalAmount = None, totalPaper = None, skip = False):
        # 电子发票路径
        self.invoicePath = './inputs'
        # images临时目录
        self.tempImagePath = './temp/images'
        # pptx临时目录
        self.tempPptxPath = './temp/pptx'
        # 凭证粘贴模板路径
        self.templatePptxPath = './凭证粘贴模板.pptx'
        # 输出路径
        self.outPath = './outputs'

        #经办人
        self.name = name
        #凭证总张数
        self.totalPage = totalPage
        #凭证总金额
        self.totalAmount = totalAmount
        #总页数
        self.totalPaper = totalPaper        
        #是否自动计算凭证总张数、总金额、总页数
        self.skip = skip

    def pyMuPDF_fitz(self, pdfPath, imagePath):
        baseName = os.path.basename(pdfPath)
        pdfName = os.path.splitext(baseName)[0]
        pdfDoc = fitz.open(pdfPath)

        page = pdfDoc[0]
        rotate = int(0)
        # 每个尺寸的缩放系数为1.3，这将为我们生成分辨率提高2.6的图像。
        # 此处若是不做设置，默认图片大小为：792X612, dpi=96
        zoom_x = 2.5 #(1.33333333-->1056x816)   (2-->1584x1224)
        zoom_y = 2.5
        mat = fitz.Matrix(zoom_x, zoom_y).preRotate(rotate)
        pix = page.getPixmap(matrix=mat, alpha=False)

        if not os.path.exists(imagePath):#判断存放图片的文件夹是否存在
            os.makedirs(imagePath) # 若图片文件夹不存在就创建

        pix.writePNG('%s/%s.png' % (imagePath, pdfName))#将图片写入指定的文件夹内

    def batchPdf2Png(self, invoicePath, outPngPath):
        files = os.listdir(invoicePath)
        pdfFiles = [f for f in files if f.endswith((".pdf"))]

        #计算总页数
        if not self.skip :
            self.totalPaper = self.totalPage = len(pdfFiles)

        #计算总金额
        if not self.skip :
            self.totalAmount = 0

        for pdfFile in pdfFiles:
            fullpath = os.path.join(invoicePath, pdfFile)
            self.pyMuPDF_fitz(fullpath, outPngPath)

            if not self.skip :
                self.totalAmount +=  float(pdfFile[:-4])


    def fillTextInSlide(self, slide, curPage, totalPage, totalAmount, curAmount, totalPaper):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text == '凭证总张数：':
                        run.text = '凭证总张数：{}张'.format(numToCN(totalPage))
                    if run.text == '本页张数：' and curAmount:
                        run.text = '本页张数：' + '壹张'
                    if run.text == '凭证总金额：':
                        run.text = '凭证总金额：¥{total:.2f}'.format(total=totalAmount)
                    if run.text == '本页金额：' and curAmount:
                        run.text = '本页金额：¥{cur:.2f}'.format(cur=curAmount)
                    if run.text == '经办人：':
                        run.text = '经办人：{}'.format(self.name)
                    if run.text == '第      页        共      页':
                        run.text = '第  {}  页        共  {}  页'.format(curPage, totalPaper)

    def batchInsertPngInSlide(self, path_to_tmpl_presentation, imgsPath):
        left, top, width, height= Cm(4.28), Cm(2.79), Cm(25.58), Cm(16.59)

        files = os.listdir(imgsPath)
        pngfiles = [f for f in files if f.endswith((".png"))]

        # 生成电子发票pptx页
        for index, pngfile in enumerate(pngfiles):
            fullpath = os.path.join(imgsPath, pngfile)

            prs = Presentation(path_to_tmpl_presentation)
            slide = prs.slides[0]
            pic = slide.shapes.add_picture(fullpath, left, top, height=height, width=width)

            # 填入文字内容
            curAmount = float(pngfile[:-4])
            self.fillTextInSlide(slide, index+1, self.totalPage, self.totalAmount, curAmount, self.totalPaper)

            pptxPath = os.path.join(self.tempPptxPath, os.path.splitext(pngfile)[0]+'.pptx')
            prs.save(pptxPath)
        
    def batchPaperInvoiceSlide(self, path_to_tmpl_presentation, imgsPath):
        files = os.listdir(imgsPath)
        # 电子发票总数
        eInvoiceCount = len(files)-1    # 忽略.gitkeep文件
        # 生成电子发票pptx页
        for index in range(eInvoiceCount+1, self.totalPaper+1):
            prs = Presentation(path_to_tmpl_presentation)
            slide = prs.slides[0]

            # 填入文字内容
            # print(index)
            self.fillTextInSlide(slide, index, self.totalPage, self.totalAmount, None, self.totalPaper)

            pptxPath = os.path.join(self.tempPptxPath, 'Page_{}.pptx'.format(index))
            prs.save(pptxPath)
        
    def init_powerpoint(self):
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        return powerpoint

    def ppt_to_pdf(self, powerpoint, inputFileName, outputFileName, formatType = 32):
        if outputFileName[-3:] != 'pdf':
            outputFileName = outputFileName.replace(".pptx","").replace(".ppt","") + ".pdf"

        deck = powerpoint.Presentations.Open(inputFileName)
        deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
        deck.Close()

    def convert_files_in_folder(self, powerpoint, folder, outPath):
        files = os.listdir(folder)
        pptfiles = [f for f in files if f.endswith((".ppt", ".pptx"))]
        for pptfile in pptfiles:
            fullpath = os.path.join(folder, pptfile)
            pdfpath = os.path.join(outPath, os.path.splitext(pptfile)[0]+'.pdf')
            self.ppt_to_pdf(powerpoint, fullpath, pdfpath)

    def del_file(self, path_data):
        for i in os.listdir(path_data) :# os.listdir(path_data)#返回一个列表，里面是当前目录下面的所有东西的相对路径
            file_data = path_data + "\\" + i#当前文件夹的下面的所有东西的绝对路径
            if os.path.isfile(file_data) == True:#os.path.isfile判断是否为文件,如果是文件,就删除.如果是文件夹.递归给del_file.
                if not file_data.endswith('.gitkeep'):
                    os.remove(file_data)
            else:
                self.del_file(file_data)

    def getfilenames(self, filepath='',filelist_out=[],file_ext='all'):
        # 遍历filepath下的所有文件，包括子目录下的文件
        for fpath, dirs, fs in os.walk(filepath):
            for f in fs:
                fi_d = os.path.join(fpath, f)
                if  file_ext == 'all':
                    filelist_out.append(fi_d)
                elif os.path.splitext(fi_d)[1] == file_ext:
                    filelist_out.append(fi_d)
                else:
                    pass
        return filelist_out

    def mergefiles(self, path, output_filename, import_bookmarks=False):
        # 遍历目录下的所有pdf将其合并输出到一个pdf文件中，输出的pdf文件默认带书签，书签名为之前的文件名
        # 默认情况下原始文件的书签不会导入，使用import_bookmarks=True可以将原文件所带的书签也导入到输出的pdf文件中
        merger = PdfFileMerger()
        filelist = self.getfilenames(filepath=path, file_ext='.pdf')

        if len(filelist) == 0:
            print("当前目录及子目录下不存在pdf文件")
            sys.exit()

        for filename in filelist:
            f = codecs.open(filename, 'rb')
            file_rd = PdfFileReader(f)
            short_filename = os.path.basename(os.path.splitext(filename)[0])
            if file_rd.isEncrypted == True:
                print('不支持的加密文件：%s'%(filename))
                continue
            merger.append(file_rd, bookmark=short_filename, import_bookmarks=import_bookmarks)
            print('合并文件：%s'%(filename))
            f.close()

        out_filename=os.path.join(os.path.abspath(path), output_filename)
        merger.write(out_filename)
        print('合并后的输出文件：%s'%(out_filename))
        merger.close()

        # # 删除单个pdf文件
        # filelist = self.getfilenames(filepath=self.outPath, file_ext='.pdf')

        # for filename in filelist:
        #     fileAbsPath = os.path.abspath(filename)
        #     print(fileAbsPath)
        #     if not fileAbsPath.endswith('All.pdf'):
        #         os.remove(fileAbsPath)

def excetue():
    if len(sys.argv) != 2 and len(sys.argv) != 5:
        print('参数个数不对')
        return

    name = str(sys.argv[1])

    totalPage = None
    totalAmount = None
    totalPaper = None
    skip = False

    if len(sys.argv) == 5:
        totalPage = int(sys.argv[2])        #凭证总张数
        totalAmount = float(sys.argv[3])    #凭证总金额
        totalPaper = int(sys.argv[4])     #总页数
        skip = True

    ic = InvoiceConverter(name, totalPage, totalAmount, totalPaper, skip)

    # pdf转图片
    ic.batchPdf2Png(ic.invoicePath, ic.tempImagePath)

    # 图片插入pptx模板
    ic.batchInsertPngInSlide(ic.templatePptxPath, ic.tempImagePath)

    # 生成纸质发票粘贴pptx页
    ic.batchPaperInvoiceSlide(ic.templatePptxPath, ic.tempImagePath)

    # pptx导出pdf
    powerpoint = ic.init_powerpoint()
    absPptxPath = os.path.abspath(ic.tempPptxPath)
    absOutPath = os.path.abspath(ic.outPath)

    ic.convert_files_in_folder(powerpoint, absPptxPath, absOutPath)

    powerpoint.Quit()

    # 清理临时文件
    ic.del_file(ic.tempImagePath)
    ic.del_file(ic.tempPptxPath)
    # 合并pdf文件
    ic.mergefiles(ic.outPath, 'All.pdf')

if __name__ == '__main__':
    excetue()